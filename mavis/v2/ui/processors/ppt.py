import pandas as pd
import io
from datetime import datetime
from pathlib import Path
from typing import List

from PIL import Image

import os
import streamlit as st

from db import ProjectDAO

class PPTExportWidget:
    def __init__(self, paths_df):
        from pptx import Presentation
        template = st.file_uploader(
            "Uploade a .ppt Template", accept_multiple_files=False,
        )
        self.prs = Presentation(template)
        self.layout = self.prs.slide_layouts[5]
        self.placeholder_id = 2
        self.slide_title = st.text_input("Slide Title", "Comparison")
        self.create_presentation(paths_df)

        path = Path(ProjectDAO().get())
        path = path / f"{datetime.now():%y%m%d}_{path.stem}.pptx"
        self.prs.save(str(path))

        with open(path, "rb") as presentation_file:
            st.download_button(
                "Download",
                presentation_file,
                mime="application/octet-stream",
                file_name=path.name
            )

    def _add_image(self, slide, heading:str, image_path:Path):

        with io.BytesIO() as bytes_like:
            Image.open(image_path).save(bytes_like, "JPEG")
            bytes_like.seek(0)
            slide.shapes.add_picture(
                bytes_like,
                10,10,
                width=slide
            )

            #slide.shapes.add_picture(img_path, pptx.util.Inches(0.5), pptx.util.Inches(1.75),
            #                 width=pptx.util.Inches(9), height=pptx.util.Inches(5))

    def create_presentation(self, images: pd.DataFrame):
        for index, row in images.iterrows():
            slide = self.prs.slides.add_slide(self.layout)
            slide.shapes.title.text = self.slide_title
            for i, (heading, path) in enumerate(zip(images.columns, row.values)):
                self._add_image(slide, heading, path)
